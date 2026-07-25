"""
Chargebotic — Investor Pitch Deck (Jul 2026)
13 slides. Two products (Kestrel charging system, Spark E self-charging drone)
plus software layer. Placeholder values render in amber brackets.

Usage:
    python3 execution/create_pitch_deck.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Palette ───────────────────────────────────────────────────────────────────
DEEP_SPACE  = RGBColor(0x0A, 0x0A, 0x0B)
CARD_BG     = RGBColor(0x12, 0x13, 0x16)
GRID_LINE   = RGBColor(0x26, 0x27, 0x2B)
BRAND       = RGBColor(0xE8, 0x49, 0x0A)
GOLD        = RGBColor(0xFF, 0xB8, 0x4D)
GREEN_CHECK = RGBColor(0x22, 0xC5, 0x5E)
WHITE       = RGBColor(0xF5, 0xF5, 0xF3)
LIGHT_GRAY  = RGBColor(0xB4, 0xB4, 0xB1)
MID_GRAY    = RGBColor(0x8A, 0x8A, 0x87)
DARK_GRAY   = RGBColor(0x5A, 0x5A, 0x57)

HEAD_FONT    = "Helvetica"
BODY_FONT    = "Helvetica"
TOTAL_SLIDES = 13
COMPANY      = "CHARGEBOTIC"
EMAIL        = "anis@chargebotic.com"
DOMAIN       = "chargebotic.com"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height


def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color=DEEP_SPACE):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def txt(slide, left, top, w, h, text, size=14, color=WHITE, bold=False,
        align=PP_ALIGN.LEFT, font=BODY_FONT):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box


def rect(slide, left, top, w, h, fill=CARD_BG, line=None, lw=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def eyebrow(slide, label, color=BRAND):
    txt(slide, Inches(0.6), Inches(0.5), Inches(11), Inches(0.3),
        label.upper(), size=10, color=color, bold=True, font=HEAD_FONT)


def title(slide, text, size=38):
    txt(slide, Inches(0.6), Inches(1.0), Inches(12.1), Inches(1.0),
        text, size=size, bold=True, font=HEAD_FONT)


def footer(slide, n):
    rect(slide, Inches(0.5), Inches(6.95), Inches(12.333), Pt(0.75), fill=GRID_LINE)
    txt(slide, Inches(0.5), Inches(7.05), Inches(3), Inches(0.3),
        COMPANY, size=8, color=MID_GRAY, bold=True, font=HEAD_FONT)
    txt(slide, Inches(4.5), Inches(7.05), Inches(4.3), Inches(0.3),
        f"{EMAIL}  ·  {DOMAIN}", size=8, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    txt(slide, Inches(10.5), Inches(7.05), Inches(2.333), Inches(0.3),
        f"{n:02d} / {TOTAL_SLIDES:02d}", size=8, color=MID_GRAY, bold=True,
        align=PP_ALIGN.RIGHT, font=HEAD_FONT)


def rows(slide, items, y0=2.2, dy=0.62, label_w=2.6, size=12):
    for i, (k, v) in enumerate(items):
        y = Inches(y0 + i * dy)
        rect(slide, Inches(0.6), y + Inches(0.1), Inches(0.05), Inches(0.34), fill=BRAND)
        txt(slide, Inches(0.85), y + Inches(0.05), Inches(label_w), Inches(0.5),
            k, size=size, color=WHITE, bold=True)
        txt(slide, Inches(0.85 + label_w + 0.2), y + Inches(0.05), Inches(11.4 - label_w), Inches(0.5),
            v, size=size, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
s = blank(); bg(s)
txt(s, Inches(0.6), Inches(0.55), Inches(9), Inches(0.3),
    "PRE-SEED  ·  JULY 2026  ·  CONFIDENTIAL", size=9, color=MID_GRAY, bold=True, font=HEAD_FONT)
txt(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(1.4),
    "CHARGEBOTIC", size=64, bold=True, font=HEAD_FONT)
txt(s, Inches(0.6), Inches(3.15), Inches(12.1), Inches(0.5),
    "The charging network for machines in the field.", size=20, color=GOLD, bold=True, font=HEAD_FONT)
txt(s, Inches(0.6), Inches(3.85), Inches(10.5), Inches(0.8),
    "Kestrel, a portable charging system that draws power from live power lines.\n"
    "Spark E, a drone that charges itself on the line. Software that meters every watt.",
    size=13, color=LIGHT_GRAY)
stats = [("$65K", "Kestrel unit price"), ("50-150W", "Harvested today (K1)"),
         ("$6M", "12-month revenue target"), ("2", "Products, one platform")]
rect(s, 0, Inches(5.35), SW, Inches(1.2), fill=RGBColor(0x0E, 0x0F, 0x12))
for i, (v, l) in enumerate(stats):
    x = Inches(0.9 + i * 3.05)
    txt(s, x, Inches(5.5), Inches(2.7), Inches(0.45), v, size=22, color=WHITE, bold=True, font=HEAD_FONT)
    txt(s, x, Inches(5.97), Inches(2.8), Inches(0.35), l, size=9, color=LIGHT_GRAY)
footer(s, 1)

# ══════════════════════════════════════════════════════════════════════════════
# 2 — PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s = blank(); bg(s)
eyebrow(s, "Problem")
title(s, "Range is the constraint on every electric machine.")
txt(s, Inches(0.6), Inches(2.0), Inches(11.8), Inches(0.6),
    "Ground: robots, cars, tractors. Air: drones, eVTOL. Useful life is capped by the battery, "
    "and every mission ends the same way: go back to base and charge.",
    size=13, color=LIGHT_GRAY)
# Tesla insight card
rect(s, Inches(0.6), Inches(3.0), Inches(11.9), Inches(1.7), fill=CARD_BG, line=GOLD, lw=1.2)
txt(s, Inches(0.95), Inches(3.25), Inches(11.2), Inches(0.4),
    "The key innovation of Tesla was not the car. It was the charging network.",
    size=18, color=WHITE, bold=True, font=HEAD_FONT)
txt(s, Inches(0.95), Inches(3.85), Inches(11.2), Inches(0.6),
    "The network made EVs valuable and useful for the first time. Machines in the field are "
    "waiting for the same unlock.", size=13, color=LIGHT_GRAY)
txt(s, Inches(0.6), Inches(5.2), Inches(11.8), Inches(0.5),
    "In the field there is no charging network. There are generators, fuel convoys, and missions cut short.",
    size=14, color=GOLD)
footer(s, 2)

# ══════════════════════════════════════════════════════════════════════════════
# 3 — SOLUTION
# ══════════════════════════════════════════════════════════════════════════════
s = blank(); bg(s)
eyebrow(s, "Solution")
title(s, "Power is everywhere. It just was never available.")
txt(s, Inches(0.6), Inches(2.0), Inches(11.5), Inches(0.6),
    "Power lines already run through nearly every operating area on earth. "
    "Kestrel makes them a charging point.",
    size=13, color=LIGHT_GRAY)
steps = [
    ("PERCH",   "Kestrel lands on the line and holds, using no flight power."),
    ("HARVEST", "It draws energy directly from the live power line."),
    ("DELIVER", "Power flows down a tether: into a battery, or direct to the machine."),
]
for i, (k, v) in enumerate(steps):
    x = Inches(0.6 + i * 4.15)
    rect(s, x, Inches(2.9), Inches(3.95), Inches(2.2), fill=CARD_BG, line=GRID_LINE)
    txt(s, x + Inches(0.25), Inches(3.15), Inches(1.2), Inches(0.5),
        f"0{i+1}", size=26, color=BRAND, bold=True, font=HEAD_FONT)
    txt(s, x + Inches(0.25), Inches(3.75), Inches(3.4), Inches(0.35),
        k, size=13, color=GOLD, bold=True, font=HEAD_FONT)
    txt(s, x + Inches(0.25), Inches(4.15), Inches(3.5), Inches(0.8),
        v, size=12, color=LIGHT_GRAY)
txt(s, Inches(0.6), Inches(5.5), Inches(11.8), Inches(0.5),
    "No fuel. No generator. No signature. A charging network that is already built.",
    size=14, color=WHITE, bold=True)
footer(s, 3)

# ══════════════════════════════════════════════════════════════════════════════
# 4 — SPARK E TEASER
# ══════════════════════════════════════════════════════════════════════════════
s = blank(); bg(s)
eyebrow(s, "And then the network charges itself")
title(s, "Spark E: drones that charge themselves on the line.")
txt(s, Inches(0.6), Inches(2.1), Inches(11.5), Inches(0.9),
    "A drone that lands on a power line, recharges its own battery, and flies on. "
    "No ground crew, no battery swaps, no return to base. Every power line on earth "
    "becomes its charging station.",
    size=14, color=LIGHT_GRAY)
stats = [("30 min", "V1 flight per charge (Orca X30)"), ("80 km/h", "V1 top speed"),
         ("120 min", "V2 flight target"), ("17,000 km", "One customer's coverage need")]
for i, (v, l) in enumerate(stats):
    x = Inches(0.6 + i * 3.1)
    rect(s, x, Inches(3.4), Inches(2.9), Inches(1.7), fill=CARD_BG, line=GRID_LINE)
    txt(s, x + Inches(0.22), Inches(3.65), Inches(2.5), Inches(0.5), v, size=24, color=GOLD, bold=True, font=HEAD_FONT)
    txt(s, x + Inches(0.22), Inches(4.25), Inches(2.5), Inches(0.7), l, size=10, color=LIGHT_GRAY)
txt(s, Inches(0.6), Inches(5.5), Inches(11.8), Inches(0.5),
    "Range extension changes the unit economics of every drone operation.",
    size=14, color=WHITE, bold=True)
footer(s, 4)

# ══════════════════════════════════════════════════════════════════════════════
# 5 — TEAM
# ══════════════════════════════════════════════════════════════════════════════
s = blank(); bg(s)
eyebrow(s, "Team")
title(s, "Built to ship hardware and win contracts.")
team = [
    ("ANIS CHERIET", "Co-Founder & CEO",
     "EV charging infrastructure across Europe. Strategy, defense partnerships, fundraising."),
    ("BO REDFEARN", "Co-Founder & CTO",
     "Formerly Apple. Designed Magline, the harvesting core. Hardware and systems."),
    ("AMIN FOROUGHI", "Engineer",
     "Computer vision and autonomy. Line detection and semi-autonomous attach."),
    ("BECKETT", "Engineer",
     "Mechanical engineering. FPV champion pilot. Flight testing."),
]
for i, (n, r, d) in enumerate(team):
    x = Inches(0.6 + i * 3.15)
    rect(s, x, Inches(2.1), Inches(2.95), Inches(2.6), fill=CARD_BG, line=GRID_LINE)
    txt(s, x + Inches(0.2), Inches(2.32), Inches(2.6), Inches(0.35), n, size=12, color=WHITE, bold=True, font=HEAD_FONT)
    txt(s, x + Inches(0.2), Inches(2.68), Inches(2.6), Inches(0.3), r, size=9, color=BRAND, bold=True, font=HEAD_FONT)
    txt(s, x + Inches(0.2), Inches(3.05), Inches(2.6), Inches(1.5), d, size=10, color=LIGHT_GRAY)
rect(s, Inches(0.6), Inches(5.0), Inches(11.9), Inches(1.1), fill=CARD_BG, line=GREEN_CHECK, lw=0.75)
txt(s, Inches(0.9), Inches(5.15), Inches(11.2), Inches(0.3),
    "ADVISORS", size=9, color=GREEN_CHECK, bold=True, font=HEAD_FONT)
txt(s, Inches(0.9), Inches(5.5), Inches(11.2), Inches(0.5),
    "Patrick Consorti (Silicon Valley operator, Bridge2)  ·  Arne Stoschek (ex Volkswagen, Better Place)  ·  "
    "[Ex-PG&E utility advisor, name TBD]",
    size=11, color=LIGHT_GRAY)
footer(s, 5)

# ══════════════════════════════════════════════════════════════════════════════
# 6 — LAUNCH PRODUCT: KESTREL
# ══════════════════════════════════════════════════════════════════════════════
s = blank(); bg(s)
eyebrow(s, "Launch product 01")
title(s, "Kestrel. Portable charging system. $65K.")
txt(s, Inches(0.6), Inches(2.0), Inches(11.5), Inches(0.5),
    "For defense: missions and vehicles operating far from any base.",
    size=13, color=LIGHT_GRAY)
flow = ["Lands on the powerline", "Tether sends power down", "Converter conditions the power", "[Optional] battery buffer"]
for i, f in enumerate(flow):
    x = Inches(0.6 + i * 3.1)
    rect(s, x, Inches(2.7), Inches(2.9), Inches(1.0), fill=CARD_BG, line=GRID_LINE)
    txt(s, x + Inches(0.2), Inches(2.85), Inches(0.6), Inches(0.4), f"{i+1}", size=18, color=BRAND, bold=True, font=HEAD_FONT)
    txt(s, x + Inches(0.7), Inches(2.88), Inches(2.1), Inches(0.75), f, size=11, color=WHITE)
rect(s, Inches(0.6), Inches(4.1), Inches(11.9), Inches(1.9), fill=CARD_BG, line=GOLD, lw=1.0)
txt(s, Inches(0.9), Inches(4.3), Inches(11), Inches(0.3), "TRACTION", size=9, color=GOLD, bold=True, font=HEAD_FONT)
txt(s, Inches(0.9), Inches(4.65), Inches(11.2), Inches(0.4),
    "Chariot Defense — battlefield power hubs, a16z-backed. Pilot in negotiation.",
    size=12, color=WHITE)
txt(s, Inches(0.9), Inches(5.15), Inches(11.2), Inches(0.4),
    "Overland — [customer details TBD]",
    size=12, color=WHITE)
footer(s, 6)

# ══════════════════════════════════════════════════════════════════════════════
# 7 — PRODUCT 2: SPARK E
# ══════════════════════════════════════════════════════════════════════════════
s = blank(); bg(s)
eyebrow(s, "Launch product 02")
title(s, "Spark E. The self-charging drone.")
cols = [
    ("V1  ·  ORCA X30", GOLD,
     ["30 min flight per charge", "Up to 80 km/h", "Payload up to [X] kg", "Autonomous modes + BVLOS"]),
    ("V2  ·  [MODEL TBD]", WHITE,
     ["120 min flight per charge", "Payload [TBD]", "Extended autonomy", "In development"]),
    ("CUSTOM", MID_GRAY,
     ["Engineering collaboration", "Customer sensors + electronics", "Integration with their autonomy", "Per-program pricing"]),
]
for i, (h, c, items) in enumerate(cols):
    x = Inches(0.6 + i * 4.15)
    rect(s, x, Inches(2.0), Inches(3.95), Inches(2.5), fill=CARD_BG, line=GRID_LINE)
    txt(s, x + Inches(0.22), Inches(2.2), Inches(3.5), Inches(0.3), h, size=11, color=c, bold=True, font=HEAD_FONT)
    for j, it in enumerate(items):
        txt(s, x + Inches(0.22), Inches(2.62 + j * 0.44), Inches(3.55), Inches(0.4),
            f"·  {it}", size=10, color=LIGHT_GRAY)
txt(s, Inches(0.6), Inches(4.85), Inches(11), Inches(0.3), "WHO BUYS IT", size=9, color=GOLD, bold=True, font=HEAD_FONT)
buyers = [
    "Cupertino — largest utility inspector in Argentina. 17,000 km of lines to cover.",
    "Automated reforestation drone fleets  ·  Aerial imagery collection  ·  Defense",
]
for i, b in enumerate(buyers):
    txt(s, Inches(0.6), Inches(5.25 + i * 0.5), Inches(11.9), Inches(0.45), b, size=12, color=WHITE)
footer(s, 7)

# ══════════════════════════════════════════════════════════════════════════════
# 8 — SOFTWARE
# ══════════════════════════════════════════════════════════════════════════════
s = blank(); bg(s)
eyebrow(s, "The layer on top")
title(s, "Software: every watt measured, every line known.")
sw = [
    ("TELEMETRY", "Live status of every Kestrel and Spark E in the field: position, power, health."),
    ("LINE HEALTH", "Every perch is an inspection. Line condition data utilities do not have today."),
    ("METERING", "Measure the energy drawn, bill for it. The mechanism that makes grid partnerships pay."),
]
for i, (k, v) in enumerate(sw):
    x = Inches(0.6 + i * 4.15)
    rect(s, x, Inches(2.3), Inches(3.95), Inches(2.4), fill=CARD_BG, line=GRID_LINE)
    txt(s, x + Inches(0.25), Inches(2.55), Inches(3.4), Inches(0.35), k, size=13, color=GOLD, bold=True, font=HEAD_FONT)
    txt(s, x + Inches(0.25), Inches(3.05), Inches(3.45), Inches(1.5), v, size=12, color=LIGHT_GRAY)
footer(s, 8)

# ══════════════════════════════════════════════════════════════════════════════
# 9 — BUSINESS MODEL
# ══════════════════════════════════════════════════════════════════════════════
s = blank(); bg(s)
eyebrow(s, "Business model")
title(s, "Hardware plus subscription.")
rows(s, [
    ("Hardware unit", "Kestrel at $65K. Spark E priced per configuration."),
    ("Software subscription", "Telemetry, line health, metering. Recurring, attached to every unit."),
    ("Training & support", "Operator training, field support, hardware upgrade path."),
], y0=2.3, dy=0.85, label_w=3.4, size=14)
footer(s, 9)

# ══════════════════════════════════════════════════════════════════════════════
# 10 — TAM
# ══════════════════════════════════════════════════════════════════════════════
s = blank(); bg(s)
eyebrow(s, "Market")
title(s, "Two markets, one charging network.")
rect(s, Inches(0.6), Inches(2.2), Inches(5.85), Inches(3.2), fill=CARD_BG, line=BRAND, lw=1.2)
txt(s, Inches(0.9), Inches(2.45), Inches(5.2), Inches(0.3), "CHARGING SYSTEMS — DEFENSE", size=11, color=BRAND, bold=True, font=HEAD_FONT)
txt(s, Inches(0.9), Inches(2.9), Inches(5.2), Inches(0.9), "250,000", size=44, color=WHITE, bold=True, font=HEAD_FONT)
txt(s, Inches(0.9), Inches(3.85), Inches(5.2), Inches(1.2),
    "Vehicles in the US fleet that operate away from fixed power. "
    "Every one is a candidate for line-drawn charging.", size=12, color=LIGHT_GRAY)
rect(s, Inches(6.85), Inches(2.2), Inches(5.85), Inches(3.2), fill=CARD_BG, line=GOLD, lw=1.2)
txt(s, Inches(7.15), Inches(2.45), Inches(5.2), Inches(0.3), "SELF-CHARGING DRONES", size=11, color=GOLD, bold=True, font=HEAD_FONT)
txt(s, Inches(7.15), Inches(2.9), Inches(5.2), Inches(0.9), "[TAM TBD]", size=44, color=GOLD, bold=True, font=HEAD_FONT)
txt(s, Inches(7.15), Inches(3.85), Inches(5.2), Inches(1.2),
    "Inspection, reforestation, imagery, defense. Every operator whose economics "
    "are capped by battery range.", size=12, color=LIGHT_GRAY)
footer(s, 10)

# ══════════════════════════════════════════════════════════════════════════════
# 11 — IS THIS LEGAL?
# ══════════════════════════════════════════════════════════════════════════════
s = blank(); bg(s)
eyebrow(s, "The question everyone asks")
title(s, "Is this legal?")
txt(s, Inches(0.6), Inches(2.0), Inches(11.5), Inches(0.5),
    "Yes, built the right way: with the utilities, not around them. Our advisor spent a career at PG&E.",
    size=14, color=LIGHT_GRAY)
rows(s, [
    ("Safety first", "Non-contact harvesting, controlled standoff, redundancies, formal safety case."),
    ("New revenue for utilities", "Metered energy drawn from their lines becomes a new revenue stream."),
    ("Line health reporting", "Every perch inspects the line. Utilities get condition data they lack today."),
], y0=2.9, dy=0.85, label_w=3.6, size=14)
txt(s, Inches(0.6), Inches(5.7), Inches(11.8), Inches(0.5),
    "Advisor: [name TBD], ex PG&E.", size=12, color=GOLD, bold=True)
footer(s, 11)

# ══════════════════════════════════════════════════════════════════════════════
# 12 — MILESTONES
# ══════════════════════════════════════════════════════════════════════════════
s = blank(); bg(s)
eyebrow(s, "Where the product is")
title(s, "Kestrel 1 works. Kestrel 2 ships Q3.", size=34)
rect(s, Inches(0.6), Inches(2.0), Inches(5.85), Inches(4.3), fill=CARD_BG, line=GREEN_CHECK, lw=1.0)
txt(s, Inches(0.9), Inches(2.2), Inches(5.2), Inches(0.3), "KESTREL 1  ·  NOW", size=11, color=GREEN_CHECK, bold=True, font=HEAD_FONT)
k1 = ["Works on powerlines, 50 to 150 W", "Demonstrated on Orca drone", "Remote piloted", "BOM $15K",
      "Pulls [X] W in [Y] minutes", "1 hr replaces a [P] generator / 300L of fuel"]
for j, it in enumerate(k1):
    txt(s, Inches(0.9), Inches(2.65 + j * 0.58), Inches(5.3), Inches(0.5), f"·  {it}", size=12, color=LIGHT_GRAY)
rect(s, Inches(6.85), Inches(2.0), Inches(5.85), Inches(4.3), fill=CARD_BG, line=GOLD, lw=1.0)
txt(s, Inches(7.15), Inches(2.2), Inches(5.2), Inches(0.3), "KESTREL 2  ·  SHIPPING Q3 2026", size=11, color=GOLD, bold=True, font=HEAD_FONT)
k2 = ["Powerlines 50 to 250 W", "Retrofit to multiple airframes", "Semi-autonomous: finds and attaches to the line",
      "Integrates with customer autonomy", "BOM $20K", "1 hr replaces a [P] generator / 900L of fuel"]
for j, it in enumerate(k2):
    txt(s, Inches(7.15), Inches(2.65 + j * 0.58), Inches(5.3), Inches(0.5), f"·  {it}", size=12, color=LIGHT_GRAY)
footer(s, 12)

# ══════════════════════════════════════════════════════════════════════════════
# 13 — NEXT 12 MONTHS
# ══════════════════════════════════════════════════════════════════════════════
s = blank(); bg(s)
eyebrow(s, "Next 12 months")
title(s, "$6M revenue. Two products in the field.")
rows(s, [
    ("Revenue", "$6M over the next 12 months."),
    ("Kestrel Gen 3", "In production at scale: 3 customers, 10 units each."),
    ("Spark E Gen 2", "30 drones sold and flying in the field."),
    ("Spark E prototype", "Proven: BOM, charge time, flight time [figures TBD]."),
    ("Non-dilutive", "SBIR Direct to Phase II."),
], y0=2.2, dy=0.72, label_w=3.2, size=14)
txt(s, Inches(0.6), Inches(6.2), Inches(11), Inches(0.35),
    "Anis Cheriet  ·  anis@chargebotic.com  ·  chargebotic.com", size=12, color=GOLD, bold=True)
footer(s, 13)

# ── Save ──────────────────────────────────────────────────────────────────────
ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
pptx_path = os.path.join(ROOT, "chargebotic-pitch-deck.pptx")
prs.save(pptx_path)
print(f"PPTX saved: {pptx_path}")

# ── Optional Google Slides upload (skipped gracefully if auth unavailable) ────
def upload_to_google_slides(local_pptx_path):
    from google.oauth2.credentials import Credentials
    from google.auth.transport.requests import Request
    from googleapiclient.discovery import build
    from googleapiclient.http import MediaFileUpload

    creds_root = os.path.dirname(os.path.dirname(ROOT))
    token_path = os.path.join(creds_root, "token.json")
    creds = None
    if os.path.exists(token_path):
        creds = Credentials.from_authorized_user_file(token_path, [
            "https://www.googleapis.com/auth/drive",
            "https://www.googleapis.com/auth/drive.file",
        ])
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            raise RuntimeError("no valid token")
    drive = build("drive", "v3", credentials=creds)
    media = MediaFileUpload(local_pptx_path, mimetype=(
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"), resumable=True)
    file = drive.files().create(body={
        "name": "Chargebotic — Investor Pitch Deck (July 2026)",
        "mimeType": "application/vnd.google-apps.presentation",
    }, media_body=media, fields="id,webViewLink").execute()
    print(f"Google Slides URL: {file.get('webViewLink')}")


try:
    upload_to_google_slides(pptx_path)
except Exception as e:
    print(f"Google Slides upload skipped ({type(e).__name__}). PPTX is ready locally.")
