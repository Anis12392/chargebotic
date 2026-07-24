"""
Chargebotic — Defense Market Research Deck
NASA/military aesthetic, 11 slides.
No external image assets required.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Palette ──────────────────────────────────────────────────────────────────
BLACK       = RGBColor(0x00, 0x00, 0x00)
DEEP_SPACE  = RGBColor(0x05, 0x08, 0x14)
ARMY_GREEN  = RGBColor(0x1A, 0x2E, 0x1A)
OD_GREEN    = RGBColor(0x3B, 0x4A, 0x2F)
AMBER       = RGBColor(0xF9, 0x73, 0x16)
DANGER_RED  = RGBColor(0xFC, 0x3D, 0x21)
STEEL_BLUE  = RGBColor(0x1A, 0x4A, 0x7A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xD4, 0xD4, 0xD4)
MID_GRAY    = RGBColor(0x9A, 0x9A, 0x9A)
DARK_GRAY   = RGBColor(0x55, 0x55, 0x55)
GRID_LINE   = RGBColor(0x22, 0x26, 0x33)
CARD_BG     = RGBColor(0x0D, 0x11, 0x1C)
GREEN_CHECK = RGBColor(0x22, 0xC5, 0x5E)

HEAD_FONT = "Helvetica"
BODY_FONT = "Helvetica"
TOTAL_SLIDES = 11

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height


# ── Helpers ──────────────────────────────────────────────────────────────────
def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color=DEEP_SPACE):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def txt(slide, left, top, w, h, text,
        size=14, color=WHITE, bold=False,
        align=PP_ALIGN.LEFT, font=BODY_FONT):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.size  = Pt(size)
    p.font.color.rgb = color
    p.font.bold  = bold
    p.font.name  = font
    p.alignment  = align
    return box


def rect(slide, left, top, w, h, fill=CARD_BG, line=None, lw=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def badge(slide, x, y, mission_id, label, color=AMBER):
    rect(slide, x, y, Inches(0.05), Inches(0.4), fill=color)
    txt(slide, x + Inches(0.18), y - Inches(0.02),
        Inches(4), Inches(0.25),
        mission_id, size=9, color=color, bold=True, font=HEAD_FONT)
    txt(slide, x + Inches(0.18), y + Inches(0.2),
        Inches(9), Inches(0.3),
        label, size=10, color=WHITE, bold=True, font=HEAD_FONT)


def footer(slide, n):
    rect(slide, Inches(0.5), Inches(6.95), Inches(12.333), Pt(0.75), fill=GRID_LINE)
    txt(slide, Inches(0.5), Inches(7.05), Inches(3), Inches(0.3),
        "CHARGEBOTIC", size=8, color=MID_GRAY, bold=True, font=HEAD_FONT)
    txt(slide, Inches(4.5), Inches(7.05), Inches(4.3), Inches(0.3),
        "anis@chargebotic.com  ·  chargebotic.com",
        size=8, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    txt(slide, Inches(10.5), Inches(7.05), Inches(2.333), Inches(0.3),
        f"{n:02d} / {TOTAL_SLIDES:02d}", size=8, color=MID_GRAY,
        bold=True, align=PP_ALIGN.RIGHT, font=HEAD_FONT)


def divider(slide, y):
    rect(slide, Inches(0.6), y, Inches(12.1), Pt(1), fill=GRID_LINE)


def stat_block(slide, x, y, w, h, stat, stat_color, label, sub=None):
    rect(slide, x, y, w, h, fill=CARD_BG, line=GRID_LINE)
    txt(slide, x + Inches(0.25), y + Inches(0.2),
        w - Inches(0.4), Inches(0.8),
        stat, size=36, color=stat_color, bold=True, font=HEAD_FONT)
    txt(slide, x + Inches(0.25), y + Inches(1.05),
        w - Inches(0.4), Inches(0.35),
        label, size=10, color=WHITE, bold=True, font=HEAD_FONT)
    if sub:
        txt(slide, x + Inches(0.25), y + Inches(1.45),
            w - Inches(0.4), Inches(0.6),
            sub, size=10, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
bg(slide, DEEP_SPACE)

# Top accent strip
rect(slide, 0, 0, SW, Inches(0.06), fill=AMBER)

badge(slide, Inches(0.6), Inches(0.5),
      "RESEARCH BRIEF — JUNE 2026", "CONFIDENTIAL · DEFENSE SEGMENT ONLY")

txt(slide, Inches(0.6), Inches(1.4), Inches(12.1), Inches(1.6),
    "DEFENSE\nMARKET RESEARCH",
    size=58, bold=True, font=HEAD_FONT)

rect(slide, Inches(0.6), Inches(3.2), Inches(3.0), Pt(2), fill=AMBER)

txt(slide, Inches(0.6), Inches(3.4), Inches(12.1), Inches(0.5),
    "TETHERED POWERLINE DRONE — BATTLEFIELD ENERGY DELIVERY",
    size=14, color=AMBER, bold=True, font=HEAD_FONT)

txt(slide, Inches(0.6), Inches(4.1), Inches(10), Inches(1.5),
    "Who buys it. How many. What they need.\nWhy now. How to reach them.",
    size=16, color=LIGHT_GRAY)

# Bottom stats bar
rect(slide, 0, Inches(6.2), SW, Inches(1.0), fill=ARMY_GREEN)
stats = [
    ("$20B+",   "DoD Annual Energy Spend"),
    ("$400/gal","Fully Burdened Fuel Cost at FOB"),
    ("59",      "US Army Brigade Combat Teams"),
    ("0",       "Direct Competitors"),
]
for i, (v, l) in enumerate(stats):
    x = Inches(1.2 + i * 3.0)
    txt(slide, x, Inches(6.25), Inches(2.0), Inches(0.4),
        v, size=22, color=WHITE, bold=True, font=HEAD_FONT)
    txt(slide, x, Inches(6.65), Inches(2.8), Inches(0.3),
        l, size=9, color=LIGHT_GRAY, font=BODY_FONT)

footer(slide, 1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
bg(slide, DEEP_SPACE)

badge(slide, Inches(0.6), Inches(0.5),
      "HAZARD — 01", "THE ENERGY LOGISTICS PROBLEM", color=DANGER_RED)

txt(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.9),
    "Fuel logistics kills soldiers.",
    size=38, bold=True, font=HEAD_FONT)

txt(slide, Inches(0.6), Inches(2.4), Inches(12), Inches(0.5),
    "Diesel generators dominate FOB power. Every gallon requires a convoy. Every convoy is a target.",
    size=13, color=LIGHT_GRAY)

# 4 pain stats
pain = [
    ("$400/gal",  DANGER_RED,  "Fully burdened fuel cost\nat a remote FOB\n(vs. $3 at the pump)"),
    ("$20B+/yr",  AMBER,       "Total DoD spend on\nfuel and electricity\nannually"),
    ("1 in 24",   DANGER_RED,  "Fuel convoys that result\nin a casualty\n(Afghanistan data)"),
    ("70–80%",    AMBER,       "Of convoy resupply weight\nthat is fuel & water —\ncrowding out ammo & food"),
]
for i, (stat, color, sub) in enumerate(pain):
    x = Inches(0.6 + i * 3.15)
    stat_block(slide, x, Inches(3.2), Inches(2.9), Inches(2.6),
               stat, color, "", sub=sub)

txt(slide, Inches(0.6), Inches(6.1), Inches(12), Inches(0.4),
    "▸  50%+ of all convoy casualties in Afghanistan in 2009 were tied to fuel resupply missions.",
    size=11, color=LIGHT_GRAY)

footer(slide, 2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — MARKET SIZE
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
bg(slide, DEEP_SPACE)

badge(slide, Inches(0.6), Inches(0.5),
      "MARKET — 02", "MARKET SIZE")

txt(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.9),
    "A $2–3B deployable power opportunity.",
    size=36, bold=True, font=HEAD_FONT)

txt(slide, Inches(0.6), Inches(2.35), Inches(12), Inches(0.4),
    "The broader military power market is $8–11B today. The relevant sub-segment (deployable off-grid tactical power) is $400M–$800M and growing fast.",
    size=12, color=LIGHT_GRAY)

# Table
table_top = Inches(3.05)
headers = ["SEGMENT", "2024 VALUE", "2030–32 PROJECTION", "CAGR"]
col_widths = [Inches(4.5), Inches(2.3), Inches(3.1), Inches(1.8)]
col_starts = [Inches(0.6)]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1] + w)

rect(slide, Inches(0.6), table_top, Inches(11.7), Inches(0.45), fill=ARMY_GREEN)
for i, (h, x, w) in enumerate(zip(headers, col_starts, col_widths)):
    txt(slide, x + Inches(0.1), table_top + Inches(0.08), w, Inches(0.3),
        h, size=10, color=WHITE, bold=True, font=HEAD_FONT)

rows = [
    ("Military Power Solutions (TAM)", "$7.7B–$10.8B",  "$13.6B–$24.9B",  "8.5–9.7%"),
    ("Military Microgrid",              "$2.0B",          "High growth",     "19.1%"),
    ("Military DC Microgrid",           "$700M",          "—",               "High"),
    ("Power Generator for Military",    "$1.2B",          "$1.7B",           "3.8%"),
    ("▸  CHARGEBOTIC SAM (deployable)", "$400M–$800M",    "$1.5B+",          "~15%"),
]
for r, (seg, val, proj, cagr) in enumerate(rows):
    row_y = table_top + Inches(0.45) + r * Inches(0.55)
    is_sam = "CHARGEBOTIC" in seg
    row_fill = RGBColor(0x12, 0x18, 0x28) if r % 2 == 0 else CARD_BG
    if is_sam:
        row_fill = RGBColor(0x1A, 0x2E, 0x1A)
    rect(slide, Inches(0.6), row_y, Inches(11.7), Inches(0.52), fill=row_fill)
    row_data = [seg, val, proj, cagr]
    for col_i, (cell, x, w) in enumerate(zip(row_data, col_starts, col_widths)):
        c = AMBER if (is_sam and col_i > 0) else (WHITE if is_sam else LIGHT_GRAY)
        bold = is_sam
        txt(slide, x + Inches(0.1), row_y + Inches(0.1), w - Inches(0.1), Inches(0.35),
            cell, size=11, color=c, bold=bold)

footer(slide, 3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — CUSTOMER 1: US ARMY
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
bg(slide, DEEP_SPACE)

badge(slide, Inches(0.6), Inches(0.5),
      "CUSTOMER — 03", "US ARMY  ·  PRIMARY BUYER")

txt(slide, Inches(0.6), Inches(1.3), Inches(8), Inches(0.9),
    "The world's largest\nFOB operator.",
    size=36, bold=True, font=HEAD_FONT)

# Left: who they are
rect(slide, Inches(0.6), Inches(2.7), Inches(6.0), Inches(3.9),
     fill=CARD_BG, line=GRID_LINE)

txt(slide, Inches(0.9), Inches(2.85), Inches(5.4), Inches(0.3),
    "WHO THEY ARE", size=9, color=AMBER, bold=True, font=HEAD_FONT)

army_facts = [
    ("59 Brigade Combat Teams", "32 active-duty + 27 National Guard"),
    ("4,000–4,700 soldiers", "Per BCT, commanded by a Colonel"),
    ("Types: 14 Infantry, 11 Armored,", "7 Stryker — each runs its own FOB"),
    ("173,000 deployed abroad", "Across 750+ US bases in 80+ countries"),
]
for i, (bold_txt, light_txt) in enumerate(army_facts):
    y = Inches(3.35 + i * 0.6)
    rect(slide, Inches(0.9), y, Inches(0.05), Inches(0.38), fill=AMBER)
    txt(slide, Inches(1.05), y + Inches(0.01), Inches(5.1), Inches(0.25),
        bold_txt, size=12, color=WHITE, bold=True)
    txt(slide, Inches(1.05), y + Inches(0.27), Inches(5.1), Inches(0.25),
        light_txt, size=10, color=LIGHT_GRAY)

# Right: pain + buying path
rect(slide, Inches(6.9), Inches(2.7), Inches(5.8), Inches(3.9),
     fill=CARD_BG, line=GRID_LINE)

txt(slide, Inches(7.2), Inches(2.85), Inches(5.2), Inches(0.3),
    "PAIN POINT", size=9, color=DANGER_RED, bold=True, font=HEAD_FONT)
txt(slide, Inches(7.2), Inches(3.15), Inches(5.2), Inches(0.8),
    "BCTs operating near power infrastructure (CONUS training, Europe, Korea) still run diesel generators — because tapping power lines at scale is not currently possible.",
    size=11, color=LIGHT_GRAY)

txt(slide, Inches(7.2), Inches(4.0), Inches(5.2), Inches(0.3),
    "HOW THEY BUY", size=9, color=AMBER, bold=True, font=HEAD_FONT)
buy_steps = [
    "Army SBIR Phase I — $250K entry point (open RFP now)",
    "Program Executive Office CS&CSS",
    "OTA (Other Transaction Authority) — faster pathway",
    "Army Futures Command (AFC) — innovation arm",
]
for i, s in enumerate(buy_steps):
    txt(slide, Inches(7.2), Inches(4.4 + i * 0.42), Inches(5.2), Inches(0.35),
        f"▸  {s}", size=11, color=WHITE)

footer(slide, 4)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — CUSTOMER 2: USSOCOM
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
bg(slide, DEEP_SPACE)

badge(slide, Inches(0.6), Inches(0.5),
      "CUSTOMER — 04", "USSOCOM  ·  HIGHEST PRIORITY BUYER", color=DANGER_RED)

txt(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.9),
    "The fastest path. The most acute need.",
    size=36, bold=True, font=HEAD_FONT)

txt(slide, Inches(0.6), Inches(2.35), Inches(12), Inches(0.4),
    "SOCOM has its own acquisition authority. They move faster, pay more, and have zero tolerance for logistics tails.",
    size=13, color=LIGHT_GRAY)

# 3 stat blocks
socom_stats = [
    ("73,000", DANGER_RED, "SOF Personnel", "Army SF, Navy SEALs, Rangers,\nMARSOC, Air Force SOF"),
    ("80+",    AMBER,      "Countries Active", "Simultaneous operations —\nzero resupply in denied zones"),
    ("Title 10\n§167", GREEN_CHECK, "Own Acquisition Authority", "Bypasses slow Army procurement;\ncan buy directly"),
]
for i, (stat, color, label, sub) in enumerate(socom_stats):
    x = Inches(0.6 + i * 4.15)
    stat_block(slide, x, Inches(3.2), Inches(3.9), Inches(2.5),
               stat, color, label, sub=sub)

# Why SOCOM first
rect(slide, Inches(0.6), Inches(5.9), Inches(11.7), Inches(0.75),
     fill=ARMY_GREEN, line=OD_GREEN)
txt(slide, Inches(0.9), Inches(6.0), Inches(11.0), Inches(0.55),
    "SOCOM SBIR: Phase I = $150K / 6 months. Phase II = up to $1M / 2 years.  ·  SOFWERX (SOF innovation hub) offers direct engagement with SOF AT&L decision makers — no prime required.",
    size=11, color=WHITE)

footer(slide, 5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — CUSTOMER 3: USMC + ALLIED FORCES
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
bg(slide, DEEP_SPACE)

badge(slide, Inches(0.6), Inches(0.5),
      "CUSTOMER — 05", "USMC + ALLIED FORCES  ·  SECONDARY BUYERS")

txt(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.9),
    "Two additional rings of buyers.",
    size=36, bold=True, font=HEAD_FONT)

# USMC card
rect(slide, Inches(0.6), Inches(2.7), Inches(5.8), Inches(3.9),
     fill=CARD_BG, line=GRID_LINE)
rect(slide, Inches(0.6), Inches(2.7), Inches(5.8), Inches(0.08), fill=STEEL_BLUE)
txt(slide, Inches(0.9), Inches(2.85), Inches(5.2), Inches(0.3),
    "US MARINE CORPS", size=10, color=STEEL_BLUE, bold=True, font=HEAD_FONT)

mc_facts = [
    ("7 Marine Expeditionary Units (MEUs)",   "Active, amphibious rapid-deployment"),
    ("~177,000 active Marines total",          ""),
    ("Expeditionary-first doctrine",           "Go ashore fast, establish power before\nsupply chains exist"),
    ("Buys via MARCORSYSCOM",                  "Often co-buys with Army on shared programs"),
    ("USMC FOB Energy program active",         "Vocal about fuel dependency reduction"),
]
for i, (b, l) in enumerate(mc_facts):
    y = Inches(3.3 + i * 0.55)
    txt(slide, Inches(0.9), y, Inches(5.0), Inches(0.25),
        b, size=11, color=WHITE, bold=True)
    if l:
        txt(slide, Inches(0.9), y + Inches(0.25), Inches(5.0), Inches(0.25),
            l, size=10, color=LIGHT_GRAY)

# Allied forces card
rect(slide, Inches(6.9), Inches(2.7), Inches(5.8), Inches(3.9),
     fill=CARD_BG, line=GRID_LINE)
rect(slide, Inches(6.9), Inches(2.7), Inches(5.8), Inches(0.08), fill=AMBER)
txt(slide, Inches(7.2), Inches(2.85), Inches(5.2), Inches(0.3),
    "NATO + PACIFIC ALLIES", size=10, color=AMBER, bold=True, font=HEAD_FONT)

allied_facts = [
    ("32 NATO member nations",         "Same doctrine, same energy problem"),
    ("Japan, South Korea, Australia",  "US ally Pacific theater — run identical FOB structures"),
    ("Foreign Military Sales (FMS)",   "US-to-ally procurement route"),
    ("Direct commercial export",       "With proper export licensing (ITAR)"),
    ("Market multiplier ~3x",          "Validated US product sells internationally"),
]
for i, (b, l) in enumerate(allied_facts):
    y = Inches(3.3 + i * 0.55)
    txt(slide, Inches(7.2), y, Inches(5.0), Inches(0.25),
        b, size=11, color=WHITE, bold=True)
    if l:
        txt(slide, Inches(7.2), y + Inches(0.25), Inches(5.0), Inches(0.25),
            l, size=10, color=LIGHT_GRAY)

footer(slide, 6)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — B2B2G STRATEGY
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
bg(slide, DEEP_SPACE)

badge(slide, Inches(0.6), Inches(0.5),
      "STRATEGY — 06", "THE B2B2G PLAY  ·  FASTEST PATH TO REVENUE", color=GREEN_CHECK)

txt(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.9),
    "Sell to the primes. They sell to the DoD.",
    size=36, bold=True, font=HEAD_FONT)

txt(slide, Inches(0.6), Inches(2.35), Inches(12), Inches(0.5),
    "Defense tech startups already hold open government contracts (IDIQs, OTAs). "
    "Chargebotic integrates as their power source. They bundle and sell to DoD under their existing vehicle.",
    size=12, color=LIGHT_GRAY)

# Three path comparison
paths = [
    ("DIRECT TO DoD", DANGER_RED, "18–36 months",
     ["SBIR Phase I → II", "Program of Record", "Long compliance overhead", "High barrier, high reward"]),
    ("B2B2G (PRIMES)", GREEN_CHECK, "3–9 months",
     ["Partner already has contract", "They handle procurement", "Faster revenue, less margin", "Proof point for direct later"]),
    ("BOTH IN PARALLEL", AMBER, "Optimal",
     ["B2B2G for speed + cash", "SBIR for long-term margin", "Diversified risk", "Recommended approach"]),
]
for i, (title, color, timeline, bullets) in enumerate(paths):
    x = Inches(0.6 + i * 4.15)
    rect(slide, x, Inches(3.1), Inches(3.9), Inches(3.5),
         fill=CARD_BG, line=color, lw=1.5)
    rect(slide, x, Inches(3.1), Inches(3.9), Inches(0.07), fill=color)
    txt(slide, x + Inches(0.25), Inches(3.25), Inches(3.4), Inches(0.3),
        title, size=10, color=color, bold=True, font=HEAD_FONT)
    txt(slide, x + Inches(0.25), Inches(3.6), Inches(3.4), Inches(0.45),
        timeline, size=28, color=WHITE, bold=True, font=HEAD_FONT)
    txt(slide, x + Inches(0.25), Inches(4.1), Inches(3.4), Inches(0.25),
        "TIME TO FIRST CONTRACT", size=8, color=MID_GRAY, font=HEAD_FONT)
    for j, b in enumerate(bullets):
        txt(slide, x + Inches(0.25), Inches(4.4 + j * 0.42), Inches(3.4), Inches(0.35),
            f"·  {b}", size=11, color=LIGHT_GRAY)

footer(slide, 7)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — B2B2G PARTNER LANDSCAPE
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
bg(slide, DEEP_SPACE)

badge(slide, Inches(0.6), Inches(0.5),
      "PARTNERS — 07", "B2B2G TARGET COMPANIES", color=GREEN_CHECK)

txt(slide, Inches(0.6), Inches(1.3), Inches(9), Inches(0.9),
    "Five companies. All need what we build.",
    size=34, bold=True, font=HEAD_FONT)

# Priority badge
rect(slide, Inches(9.5), Inches(1.25), Inches(3.2), Inches(0.6),
     fill=RGBColor(0x0D, 0x1F, 0x0D), line=GREEN_CHECK, lw=1.0)
txt(slide, Inches(9.6), Inches(1.35), Inches(3.0), Inches(0.4),
    "★  CALL CHARIOT FIRST", size=11, color=GREEN_CHECK, bold=True, font=HEAD_FONT)

# Partner table
table_top = Inches(2.55)
headers = ["COMPANY", "FUNDING", "WHAT THEY DO", "WHY THEY NEED CHARGEBOTIC", "PRIORITY"]
col_widths = [Inches(2.0), Inches(1.5), Inches(2.8), Inches(4.0), Inches(1.3)]
col_starts = [Inches(0.6)]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1] + w)

rect(slide, Inches(0.6), table_top, Inches(11.7), Inches(0.42), fill=ARMY_GREEN)
for h, x, w in zip(headers, col_starts, col_widths):
    txt(slide, x + Inches(0.08), table_top + Inches(0.08), w, Inches(0.28),
        h, size=9, color=WHITE, bold=True, font=HEAD_FONT)

partners = [
    ("Chariot Defense",   "$41M · a16z",
     "Software-defined power distribution (Amphora)",
     "Routes power but has NO source — Chargebotic IS the source",
     "★★★", GREEN_CHECK),
    ("Anduril",           "$20B+ val.",
     "Autonomous drones, sensors, directed energy, Lattice OS",
     "All their edge systems are power-hungry — we supply the line",
     "★★★", AMBER),
    ("Shield AI",         "Pre-IPO",
     "Hivemind autonomy software + V-BAT drone platform",
     "Drones need persistent power in comms-denied environments",
     "★★☆", AMBER),
    ("Mach Industries",   "$1.8B val.",
     "UAV manufacturing at scale for DoD",
     "Power logistics is their customers' #1 pain — add-on hardware",
     "★★☆", MID_GRAY),
    ("Windlift",          "$24M DoD",
     "Tethered drone generating power from wind",
     "Same delivery model, different source — partner or watch closely",
     "★☆☆", MID_GRAY),
]
for r, (name, funding, what, why, priority, color) in enumerate(partners):
    row_y = table_top + Inches(0.42) + r * Inches(0.72)
    row_fill = RGBColor(0x0D, 0x1F, 0x0D) if r == 0 else (
        RGBColor(0x12, 0x18, 0x28) if r % 2 == 0 else CARD_BG)
    rect(slide, Inches(0.6), row_y, Inches(11.7), Inches(0.69), fill=row_fill)
    cells = [name, funding, what, why, priority]
    for cell, x, w in zip(cells, col_starts, col_widths):
        c = color if cell == name else (GREEN_CHECK if cell == priority and r == 0 else LIGHT_GRAY)
        bold = (cell == name)
        txt(slide, x + Inches(0.08), row_y + Inches(0.17),
            w - Inches(0.1), Inches(0.4),
            cell, size=10, color=c, bold=bold)

# Bottom warning
rect(slide, Inches(0.6), Inches(6.55), Inches(11.7), Inches(0.55),
     fill=RGBColor(0x1A, 0x10, 0x05), line=AMBER, lw=0.75)
txt(slide, Inches(0.9), Inches(6.65), Inches(11.0), Inches(0.35),
    "⚠  IP PROTECTION REQUIRED:  File patents before any demo conversation. NDA + term sheet fast. Never demo before IP is secured.",
    size=11, color=AMBER, bold=True)

footer(slide, 8)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — COMPETITIVE LANDSCAPE
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
bg(slide, DEEP_SPACE)

badge(slide, Inches(0.6), Inches(0.5),
      "LANDSCAPE — 08", "COMPETITIVE ANALYSIS")

txt(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.9),
    "No direct competitor. White space.",
    size=36, bold=True, font=HEAD_FONT)

txt(slide, Inches(0.6), Inches(2.35), Inches(12), Inches(0.4),
    "No product today autonomously taps a power line and delivers that energy via tether to a ground station.",
    size=13, color=LIGHT_GRAY)

# Table
table_top = Inches(3.05)
headers = ["PLAYER", "WHAT THEY DO", "FUNDING", "DIFFERENCE FROM CHARGEBOTIC"]
col_widths = [Inches(2.3), Inches(3.2), Inches(1.8), Inches(4.4)]
col_starts = [Inches(0.6)]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1] + w)

rect(slide, Inches(0.6), table_top, Inches(11.7), Inches(0.45), fill=ARMY_GREEN)
for h, x, w in zip(headers, col_starts, col_widths):
    txt(slide, x + Inches(0.1), table_top + Inches(0.08), w, Inches(0.3),
        h, size=9, color=WHITE, bold=True, font=HEAD_FONT)

rows = [
    ("CHARGEBOTIC",    "Power line tap → tether → ground battery",
     "—",              "✓  UNIQUE — NO EQUIVALENT EXISTS", GREEN_CHECK),
    ("Windlift",       "Tethered drone powered by wind",
     "$24M DoD",       "Wind-dependent, no infrastructure tap", MID_GRAY),
    ("DARPA POWER",    "Laser-based wireless power relay (RTX, Draper)",
     "$10M",           "Optical/wireless, different physics, no tether", MID_GRAY),
    ("Solar + Battery","Deployed solar panels + battery banks",
     "Widespread",     "Weather-dependent, 30-50% fuel reduction only", MID_GRAY),
    ("Diesel Generators","Standard Army MEP generator sets",
     "Incumbent",      "Fuel-dependent, loud, logistics tail", MID_GRAY),
]
for r, row in enumerate(rows):
    if len(row) == 5:
        player, what, funding, diff, diff_color = row
    else:
        player, what, funding, diff = row
        diff_color = LIGHT_GRAY
    row_y = table_top + Inches(0.45) + r * Inches(0.57)
    is_us = r == 0
    row_fill = RGBColor(0x12, 0x18, 0x28) if r % 2 == 0 else CARD_BG
    if is_us:
        row_fill = RGBColor(0x0D, 0x1F, 0x0D)
    rect(slide, Inches(0.6), row_y, Inches(11.7), Inches(0.54), fill=row_fill)
    cells = [player, what, funding, diff]
    cell_colors = [
        AMBER if is_us else WHITE,
        LIGHT_GRAY, LIGHT_GRAY,
        diff_color if is_us else MID_GRAY
    ]
    for cell, x, w, c in zip(cells, col_starts, col_widths, cell_colors):
        txt(slide, x + Inches(0.1), row_y + Inches(0.1), w - Inches(0.1), Inches(0.4),
            cell, size=10, color=c, bold=(is_us and cells.index(cell) == 0))

footer(slide, 9)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — KEY PROGRAMS & ENTRY PATH
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
bg(slide, DEEP_SPACE)

badge(slide, Inches(0.6), Inches(0.5),
      "EXECUTION — 09", "KEY PROGRAMS & PROCUREMENT PATH")

txt(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.9),
    "How to get in the door.",
    size=36, bold=True, font=HEAD_FONT)

# Left: active programs (buying signals)
rect(slide, Inches(0.6), Inches(2.7), Inches(5.8), Inches(4.0),
     fill=CARD_BG, line=GRID_LINE)
txt(slide, Inches(0.9), Inches(2.85), Inches(5.2), Inches(0.3),
    "ACTIVE DoD PROGRAMS (BUYING SIGNALS)", size=9, color=AMBER, bold=True, font=HEAD_FONT)

programs = [
    ("Army SBIR Power Solutions",
     "$250K Phase I — open RFP specifically for military power innovation"),
    ("SOCOM SBIR",
     "$150K Phase I / $1M Phase II — SOF power & autonomy"),
    ("DARPA POWER",
     "$10M+ for novel energy delivery — signals DoD appetite"),
    ("Army Operational Energy Strategy",
     "Formal mandate to reduce FOB fuel dependency"),
    ("Tethered Drones → Army Program",
     "Army making tethered drones formal in 2025 — direct precedent"),
    ("DoD/GAO Convoy Reports",
     "Congressional pressure = political cover for procurement"),
]
for i, (p, desc) in enumerate(programs):
    y = Inches(3.25 + i * 0.52)
    rect(slide, Inches(0.9), y + Inches(0.05), Inches(0.05), Inches(0.3), fill=AMBER)
    txt(slide, Inches(1.05), y + Inches(0.03), Inches(5.0), Inches(0.22),
        p, size=11, color=WHITE, bold=True)
    txt(slide, Inches(1.05), y + Inches(0.26), Inches(5.0), Inches(0.22),
        desc, size=9, color=LIGHT_GRAY)

# Right: entry sequence
rect(slide, Inches(6.9), Inches(2.7), Inches(5.8), Inches(4.0),
     fill=CARD_BG, line=GRID_LINE)
txt(slide, Inches(7.2), Inches(2.85), Inches(5.2), Inches(0.3),
    "RECOMMENDED ENTRY SEQUENCE", size=9, color=GREEN_CHECK, bold=True, font=HEAD_FONT)

steps = [
    ("01", AMBER,       "Army SBIR Phase I",
     "$250K · 6 months · lowest barrier\nArmy has open power opportunity now"),
    ("02", AMBER,       "SOCOM SBIR Phase I",
     "$150K · run in parallel · faster approval\nDirect access via SOFWERX"),
    ("03", STEEL_BLUE,  "OTA Prototype Agreement",
     "Post-Phase II · no competitive bid required\nBridge to production contract"),
    ("04", GREEN_CHECK, "Program of Record",
     "Long-term goal · becomes budget line item\nin Army or SOCOM procurement"),
]
for i, (num, color, title, desc) in enumerate(steps):
    y = Inches(3.25 + i * 0.82)
    txt(slide, Inches(7.2), y, Inches(0.6), Inches(0.5),
        num, size=24, color=color, bold=True, font=HEAD_FONT)
    txt(slide, Inches(7.9), y, Inches(4.5), Inches(0.28),
        title, size=12, color=WHITE, bold=True)
    txt(slide, Inches(7.9), y + Inches(0.3), Inches(4.5), Inches(0.45),
        desc, size=10, color=LIGHT_GRAY)

footer(slide, 10)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — SUMMARY
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
bg(slide, DEEP_SPACE)

rect(slide, 0, 0, SW, Inches(0.06), fill=AMBER)

badge(slide, Inches(0.6), Inches(0.5),
      "SUMMARY — 08", "THE FULL PICTURE")

txt(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.9),
    "Large market. No competitor. Two paths to revenue.",
    size=32, bold=True, font=HEAD_FONT)

# Two-column summary table
left_items = [
    ("TAM",               "$7.7B–$10.8B",   "Military power market, 2024"),
    ("SAM",               "$400M–$800M",     "Deployable off-grid tactical power"),
    ("Market growth",     "8.5–19% CAGR",   "Depending on sub-segment"),
    ("Direct customer 1", "US Army",         "32 active BCTs, ~150K deployed soldiers"),
    ("Direct customer 2", "USSOCOM",         "73K SOF, own acquisition authority, fastest path"),
]
right_items = [
    ("B2B2G target #1",  "Chariot Defense", "$41M a16z-backed — routes power, needs our source"),
    ("B2B2G target #2",  "Anduril",         "$20B+ val — edge systems starved of power"),
    ("Direct competitors","None",            "White space — no equivalent product exists"),
    ("Fuel at FOBs",      "$400/gallon",     "Fully burdened — extreme willingness to pay"),
    ("First move",        "IP + B2B2G",      "File patents → call Chariot → run SBIR in parallel"),
]

for col_i, items in enumerate([left_items, right_items]):
    x_base = Inches(0.6 + col_i * 6.4)
    rect(slide, x_base, Inches(2.7), Inches(6.0), Inches(3.85),
         fill=CARD_BG, line=GRID_LINE)
    for r, (label, value, note) in enumerate(items):
        row_y = Inches(2.85 + r * 0.68)
        if r > 0:
            rect(slide, x_base + Inches(0.2), row_y - Inches(0.08),
                 Inches(5.6), Pt(0.5), fill=GRID_LINE)
        txt(slide, x_base + Inches(0.25), row_y + Inches(0.02),
            Inches(1.8), Inches(0.28),
            label, size=9, color=MID_GRAY, bold=True, font=HEAD_FONT)
        txt(slide, x_base + Inches(2.1), row_y + Inches(0.02),
            Inches(1.8), Inches(0.28),
            value, size=13, color=AMBER, bold=True, font=HEAD_FONT)
        txt(slide, x_base + Inches(0.25), row_y + Inches(0.32),
            Inches(5.5), Inches(0.28),
            note, size=10, color=LIGHT_GRAY)

# Bottom conclusion bar
rect(slide, 0, Inches(6.75), SW, Inches(0.6), fill=ARMY_GREEN)
txt(slide, Inches(0.6), Inches(6.82), Inches(12.1), Inches(0.4),
    "BOTTOM LINE:  $20B+ DoD energy problem. No direct competitor. Call Chariot Defense first (they distribute power but have no source — we are the source). File IP. Run B2B2G and SBIR in parallel.",
    size=11, color=WHITE, font=HEAD_FONT)

footer(slide, 11)


# ── Save ─────────────────────────────────────────────────────────────────────
output = os.path.join(os.path.dirname(os.path.dirname(__file__)),
                      "defense-market-research.pptx")
prs.save(output)
print(f"Saved: {output}")
